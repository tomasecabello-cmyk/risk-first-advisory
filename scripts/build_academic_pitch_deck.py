#!/usr/bin/env python3
"""
Build Risk-First Advisory ACADEMIC pitch deck (12 slides, 16:9 widescreen).

Tono: académico-financiero / institucional / wealth-management.
Audiencia: profesor de finanzas, director de carrera, asesor o evaluador
no técnico.

Distribución del contenido:
   60 % marco financiero, problema y metodología
   25 % producto / demo
   15 % estado actual, limitaciones, roadmap

NO reemplaza al pitch deck técnico previo (Risk-First-Advisory-Pitch-Deck.pptx).
Sale en archivo separado: Risk-First-Advisory-Academic-Deck.pptx

Run: python scripts/build_academic_pitch_deck.py
"""

from __future__ import annotations

import math

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ─────────────────────────────────────────────────────────────────────
# Palette  (misma del producto y del pitch deck — coherencia visual)
# ─────────────────────────────────────────────────────────────────────
NAVY_DARK  = RGBColor(0x0A, 0x1A, 0x36)
NAVY_900   = RGBColor(0x0B, 0x1E, 0x3F)
NAVY_800   = RGBColor(0x12, 0x2A, 0x4F)
NAVY_700   = RGBColor(0x1E, 0x3A, 0x6B)
NAVY_600   = RGBColor(0x1D, 0x4E, 0xD8)
NAVY_500   = RGBColor(0x25, 0x63, 0xEB)
CYAN_500   = RGBColor(0x0E, 0xA5, 0xE9)
CYAN_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
SKY_LIGHT  = RGBColor(0xBF, 0xDB, 0xFE)
EMERALD    = RGBColor(0x05, 0x96, 0x69)
EMERALD_LT = RGBColor(0xD1, 0xFA, 0xE5)
AMBER      = RGBColor(0xD9, 0x77, 0x06)
AMBER_LT   = RGBColor(0xFE, 0xF3, 0xC7)
ROSE       = RGBColor(0xDC, 0x26, 0x26)
VIOLET     = RGBColor(0x6D, 0x28, 0xD9)
VIOLET_LT  = RGBColor(0xED, 0xE9, 0xFE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF8, 0xFA, 0xFC)
BG_LIGHT   = RGBColor(0xF7, 0xF8, 0xFC)   # very light, more academic feel
BG_PAPER   = RGBColor(0xFA, 0xFB, 0xFD)
BG_SOFT    = RGBColor(0xEE, 0xF1, 0xF7)
BORDER     = RGBColor(0xE5, 0xE7, 0xEB)
BORDER_STR = RGBColor(0xCB, 0xD5, 0xE1)
TEXT_DARK  = RGBColor(0x0F, 0x17, 0x2A)
TEXT_BODY  = RGBColor(0x33, 0x41, 0x55)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
TEXT_FAINT = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_INK  = RGBColor(0xE2, 0xE8, 0xF0)
SOFT_INK   = RGBColor(0xCB, 0xD5, 0xE1)

# Typography — Cambria para cabezas (académico), Calibri para body
FONT_HEAD  = "Cambria"
FONT_BODY  = "Calibri"
FONT_QUOTE = "Cambria"
FONT_MONO  = "Consolas"


# ─────────────────────────────────────────────────────────────────────
# Presentation setup (16:9 widescreen)
# ─────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
prs.core_properties.title  = "Risk-First Advisory — Academic deck"
prs.core_properties.author = "Risk-First Advisory"
BLANK = prs.slide_layouts[6]

SW = 13.333
SH = 7.5


# ─────────────────────────────────────────────────────────────────────
# Helpers (mismas primitivas que el pitch deck, self-contained)
# ─────────────────────────────────────────────────────────────────────

def add_slide(bg=BG_LIGHT):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                           prs.slide_width, prs.slide_height)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def rect(slide, x, y, w, h, *, fill=WHITE, line=None, line_w=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is not None:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def round_rect(slide, x, y, w, h, *, fill=WHITE, line=None, line_w=0.5, radius=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    s.adjustments[0] = radius
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is not None:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def oval(slide, x, y, d, *, fill, line=None, h=None):
    H = h if h is not None else d
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(x), Inches(y), Inches(d), Inches(H))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is not None:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def hline(slide, x1, y1, x2, y2, *, color=BORDER, weight=0.5):
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1),
                                    Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def arrow(slide, x1, y1, x2, y2, *, color=TEXT_FAINT, weight=1.0):
    """Thin horizontal arrow connector."""
    ln = slide.shapes.add_connector(2, Inches(x1), Inches(y1),
                                    Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def _set_char_spacing(run, hundredths_pt: int):
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(hundredths_pt)))


def text(slide, x, y, w, h, body, *, size=14, color=TEXT_BODY, bold=False,
         italic=False, align="left", valign="top", font=FONT_BODY, charspace=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    if isinstance(body, list):
        for seg in body:
            if isinstance(seg, str):
                seg_text, seg_opts = seg, {}
            else:
                seg_text, seg_opts = seg
            r_ = p.add_run()
            r_.text = seg_text
            f = r_.font
            f.name   = seg_opts.get("font", font)
            f.size   = Pt(seg_opts.get("size", size))
            f.bold   = seg_opts.get("bold", bold)
            f.italic = seg_opts.get("italic", italic)
            f.color.rgb = seg_opts.get("color", color)
            cs = seg_opts.get("charspace", charspace)
            if cs:
                _set_char_spacing(r_, cs)
    else:
        r_ = p.add_run()
        r_.text = body
        f = r_.font
        f.name   = font
        f.size   = Pt(size)
        f.bold   = bold
        f.italic = italic
        f.color.rgb = color
        if charspace:
            _set_char_spacing(r_, charspace)
    return tb


def text_lines(slide, x, y, w, h, lines, *, size=14, color=TEXT_BODY, bold=False,
               align="left", valign="top", font=FONT_BODY, line_spacing=1.20,
               space_after=4):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}[align]
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        segs = ln if isinstance(ln, list) else [ln]
        for seg in segs:
            if isinstance(seg, str):
                seg_text, seg_opts = seg, {}
            else:
                seg_text, seg_opts = seg
            r_ = p.add_run()
            r_.text = seg_text
            f = r_.font
            f.name   = seg_opts.get("font", font)
            f.size   = Pt(seg_opts.get("size", size))
            f.bold   = seg_opts.get("bold", bold)
            f.italic = seg_opts.get("italic", False)
            f.color.rgb = seg_opts.get("color", color)
    return tb


def eyebrow_chip(slide, x, y, label, *, fg=WHITE, bg=NAVY_900, w=2.3, h=0.34):
    round_rect(slide, x, y, w, h, fill=bg, radius=0.45)
    text(slide, x, y, w, h, label, size=10.5, color=fg, bold=True,
         align="center", valign="middle", charspace=160)


def card(slide, x, y, w, h, *, fill=WHITE, accent=None, accent_w=0.07,
         line=BORDER, radius=0.06):
    round_rect(slide, x, y, w, h, fill=fill, line=line, radius=radius)
    if accent is not None:
        rect(slide, x, y, accent_w, h, fill=accent)


def numbered_circle(slide, x, y, d, num, *, bg=NAVY_900, fg=WHITE, size=14):
    oval(slide, x, y, d, fill=bg)
    text(slide, x, y, d, d, str(num), size=size, color=fg, bold=True,
         align="center", valign="middle")


def slide_title_block(slide, eyebrow, title, sub=None, *,
                      eyebrow_bg=NAVY_900, eyebrow_w=2.5, sub_color=TEXT_MUTED):
    """Estandar para cada slide de contenido (light bg)."""
    eyebrow_chip(slide, 0.85, 0.55, eyebrow, fg=WHITE, bg=eyebrow_bg, w=eyebrow_w)
    text(slide, 0.85, 0.98, 11.6, 0.85, title,
         size=28, color=NAVY_900, bold=True, font=FONT_HEAD)
    if sub:
        text(slide, 0.85, 1.78, 11.6, 0.50, sub,
             size=14, color=sub_color, italic=True, font=FONT_BODY)


def footer(slide, num, total=12):
    text(slide, 0.55, SH - 0.4, 5.0, 0.3,
         "Risk-First Advisory  ·  Academic deck",
         size=9, color=TEXT_FAINT, align="left", font=FONT_BODY)
    text(slide, SW - 1.3, SH - 0.4, 1.0, 0.3,
         f"{num} / {total}", size=9, color=TEXT_FAINT, align="right",
         font=FONT_BODY)


# ─────────────────────────────────────────────────────────────────────
# Slide 1 — Portada
# ─────────────────────────────────────────────────────────────────────
def slide_1_portada():
    s = add_slide(bg=NAVY_DARK)

    # Decoración mínima: 2 líneas verticales sutiles a la derecha (no acent bars)
    for x in [11.4, 12.0]:
        hline(s, x, 0.0, x, SH, color=NAVY_700, weight=0.5)

    # Brand mark (top-left)
    oval(s, 0.85, 0.78, 0.18, fill=CYAN_500)
    text(s, 1.15, 0.66, 6.0, 0.5,
         "RISK-FIRST  ADVISORY", size=11, color=SKY_LIGHT, bold=True,
         charspace=240, font=FONT_BODY)

    # Eyebrow chip arriba del título
    eyebrow_chip(s, 0.85, 1.85,
                 "PROYECTO DE FINANZAS APLICADO CON IA  ·  ACADEMIC DECK",
                 fg=CYAN_LIGHT, bg=NAVY_800, w=6.4)

    # Título grande + subtítulo
    text(s, 0.85, 2.40, 11.6, 1.20,
         "Risk-First Advisory", size=62, bold=True, color=WHITE,
         font=FONT_HEAD)
    text_lines(s, 0.85, 3.75, 11.6, 1.40, [
        [("IA supervisada por asesores para convertir perfilamiento,",
          {"size": 22, "color": LIGHT_INK, "font": FONT_BODY})],
        [("suitability y construcción de carteras",
          {"size": 22, "color": LIGHT_INK, "font": FONT_BODY})],
        [("en un proceso justificable y auditable.",
          {"size": 22, "color": LIGHT_INK, "font": FONT_BODY})],
    ], line_spacing=1.20, space_after=2)

    # Frase fuerte central
    text(s, 0.85, 5.55, 11.6, 0.65,
         "“La IA propone.  El asesor decide.”",
         size=22, color=CYAN_500, italic=True, font=FONT_QUOTE)

    # Línea horizontal sutil de cierre
    hline(s, 0.85, 6.30, 6.5, 6.30, color=NAVY_700, weight=1.0)

    # Autoría / contexto académico
    text(s, 0.85, 6.40, 11.6, 0.35,
         "Marco teórico · metodología · prototipo funcional",
         size=12.5, color=SOFT_INK, font=FONT_BODY)

    # Disclaimer pequeño abajo a la derecha
    text(s, SW - 6.4, SH - 0.55, 6.0, 0.35,
         "Prototipo funcional local · no production-ready",
         size=10, color=TEXT_FAINT, italic=True, align="right",
         font=FONT_BODY)


# ─────────────────────────────────────────────────────────────────────
# Slide 2 — Problema financiero
# ─────────────────────────────────────────────────────────────────────
def slide_2_problema():
    s = add_slide(bg=BG_LIGHT)

    slide_title_block(s,
        "EL PROBLEMA",
        "Recomendar carteras exige más que elegir instrumentos",
        sub="La decisión de inversión debería partir del cliente, "
            "no del producto — y debería poder reconstruirse después.",
        eyebrow_w=1.7)

    # 6 puntos en grid 3x2 — texto breve, sin numbering pesado
    items = [
        ("El asesor primero debe entender al cliente",
         "Sin un perfil bien construido, cualquier recomendación es opcional."),
        ("El KYC suele quedar fragmentado",
         "Forms, planillas, conversaciones — sin una vista única ni versionada."),
        ("El perfil de riesgo puede ser subjetivo",
         "Dos asesores arman perfiles distintos para el mismo cliente."),
        ("Las recomendaciones no son fácilmente trazables",
         "El cliente recibe una cartera, pero el camino de la decisión se pierde."),
        ("Los overrides quedan mal documentados",
         "Cuando se excede el presupuesto de riesgo, la justificación rara vez es formal."),
        ("La IA puede amplificar estos problemas",
         "Si no está supervisada y registrada, traslada los huecos a otra capa."),
    ]
    cols = 3
    pad_x = 0.85
    pad_top = 2.55
    gap_x = 0.30
    gap_y = 0.30
    cw = (SW - 2 * pad_x - (cols - 1) * gap_x) / cols
    ch = 1.85

    for i, (title, body) in enumerate(items):
        col = i % cols
        row = i // cols
        x = pad_x + col * (cw + gap_x)
        y = pad_top + row * (ch + gap_y)
        card(s, x, y, cw, ch, fill=WHITE, accent=NAVY_600, accent_w=0.05)
        text(s, x + 0.30, y + 0.20, cw - 0.45, 0.62, title,
             size=14, color=NAVY_900, bold=True, font=FONT_HEAD)
        text(s, x + 0.30, y + 0.90, cw - 0.45, 0.85, body,
             size=11, color=TEXT_BODY, font=FONT_BODY)

    footer(s, 2)


# ─────────────────────────────────────────────────────────────────────
# Slide 3 — Marco teórico: suitability
# ─────────────────────────────────────────────────────────────────────
def slide_3_suitability():
    s = add_slide(bg=BG_LIGHT)

    slide_title_block(s,
        "MARCO TEÓRICO  ·  I",
        "Suitability antes que producto",
        sub="Una recomendación de inversión debería partir del cliente "
            "y de los condicionantes que su situación impone.",
        eyebrow_w=3.0)

    # Quote central
    text(s, 0.85, 2.45, 11.6, 0.60,
         "“Antes de hablar de cartera, hay que saber qué cartera "
         "es adecuada para ese inversor.”",
         size=18, italic=True, color=NAVY_900, font=FONT_QUOTE)

    # 7 dimensiones del perfil — 4 + 3 layout
    dims = [
        ("01", "Objetivos del inversor",   "Qué busca lograr y en qué plazo."),
        ("02", "Horizonte temporal",       "Cuándo necesita acceder al capital."),
        ("03", "Tolerancia al riesgo",     "Qué nivel de pérdida tolera emocionalmente."),
        ("04", "Capacidad de riesgo",      "Qué nivel de pérdida puede absorber sin dañar su plan."),
        ("05", "Necesidad de liquidez",    "Qué porcentaje debería estar disponible en el corto plazo."),
        ("06", "Experiencia previa",       "Familiaridad con instrumentos y volatilidad."),
        ("07", "Restricciones y preferencias", "ESG, jurisdicción, monedas, exclusiones, productos prohibidos."),
    ]
    cols = 4
    pad_x = 0.85
    cy0 = 3.40
    cw = (SW - 2 * pad_x - (cols - 1) * 0.20) / cols
    ch = 1.55

    for i, (n, title, body) in enumerate(dims):
        col = i % cols
        row = i // cols
        x = pad_x + col * (cw + 0.20)
        y = cy0 + row * (ch + 0.20)
        round_rect(s, x, y, cw, ch, fill=WHITE, line=BORDER, radius=0.08)
        text(s, x + 0.30, y + 0.18, cw - 0.45, 0.35,
             n, size=11, color=CYAN_500, bold=True, font=FONT_MONO,
             charspace=160)
        text(s, x + 0.30, y + 0.50, cw - 0.45, 0.45,
             title, size=12.5, color=NAVY_900, bold=True, font=FONT_HEAD)
        text(s, x + 0.30, y + 0.92, cw - 0.45, 0.55,
             body, size=10, color=TEXT_BODY, font=FONT_BODY)

    footer(s, 3)


# ─────────────────────────────────────────────────────────────────────
# Slide 4 — Marco teórico: risk-first vs return-chasing
# ─────────────────────────────────────────────────────────────────────
def slide_4_risk_first():
    s = add_slide(bg=BG_LIGHT)

    slide_title_block(s,
        "MARCO TEÓRICO  ·  II",
        "Risk-first: empezar por el riesgo aceptable, no por el retorno deseado",
        sub="Cambiar el orden mental del proceso cambia la calidad "
            "(y la auditabilidad) de la decisión.",
        eyebrow_w=3.0)

    # Two-column comparison
    col_w = 5.7
    col_h = 3.65
    cy = 2.55
    gap = 0.30
    left_x = 0.85
    right_x = left_x + col_w + gap

    # LEFT — Return-chasing
    card(s, left_x, cy, col_w, col_h, fill=WHITE, accent=ROSE, accent_w=0.06)
    text(s, left_x + 0.35, cy + 0.30, col_w - 0.5, 0.35,
         "ENFOQUE TRADICIONAL  ·  RETURN-CHASING",
         size=10.5, color=ROSE, bold=True, charspace=160, font=FONT_BODY)
    text(s, left_x + 0.35, cy + 0.80, col_w - 0.5, 0.50,
         "Empieza por el retorno deseado", size=18,
         color=NAVY_900, bold=True, font=FONT_HEAD)
    bullets_left = [
        "El cliente pide cierto número (“10 % anual”).",
        "Se buscan productos que “pinten” ese número.",
        "El riesgo aparece como consecuencia.",
        "Exceder límites se discute, no se firma.",
        "La justificación llega después, no antes.",
    ]
    for i, b in enumerate(bullets_left):
        y_ = cy + 1.50 + i * 0.38
        oval(s, left_x + 0.45, y_ + 0.13, 0.10, fill=ROSE)
        text(s, left_x + 0.70, y_, col_w - 0.85, 0.40,
             b, size=12, color=TEXT_BODY, font=FONT_BODY)

    # RIGHT — Risk-first
    card(s, right_x, cy, col_w, col_h, fill=WHITE, accent=EMERALD, accent_w=0.06)
    text(s, right_x + 0.35, cy + 0.30, col_w - 0.5, 0.35,
         "RISK-FIRST",
         size=10.5, color=EMERALD, bold=True, charspace=160, font=FONT_BODY)
    text(s, right_x + 0.35, cy + 0.80, col_w - 0.5, 0.50,
         "Empieza por el riesgo aceptable", size=18,
         color=NAVY_900, bold=True, font=FONT_HEAD)
    bullets_right = [
        "El KYC fija el presupuesto de riesgo aprobado.",
        "El retorno se busca dentro de ese límite.",
        "Exceder el límite exige override firmado.",
        "Cada paso queda atado a la justificación.",
        "La decisión es reconstruible ex-post.",
    ]
    for i, b in enumerate(bullets_right):
        y_ = cy + 1.50 + i * 0.38
        oval(s, right_x + 0.45, y_ + 0.13, 0.10, fill=EMERALD)
        text(s, right_x + 0.70, y_, col_w - 0.85, 0.40,
             b, size=12, color=TEXT_BODY, font=FONT_BODY)

    # Frase fuerte abajo
    text(s, 0.85, 6.40, 11.6, 0.55,
         "“No se optimiza contra el deseo de retorno; "
         "se optimiza dentro del riesgo aprobado.”",
         size=15, color=NAVY_900, italic=True, align="center", font=FONT_QUOTE)

    footer(s, 4)


# ─────────────────────────────────────────────────────────────────────
# Slide 5 — Marco teórico: carteras bajo restricciones
# ─────────────────────────────────────────────────────────────────────
def slide_5_carteras_restricciones():
    s = add_slide(bg=BG_LIGHT)

    slide_title_block(s,
        "MARCO TEÓRICO  ·  III",
        "De la teoría de carteras a una decisión asesorada",
        sub="No todas las carteras eficientes son adecuadas para todos los clientes. "
            "El perfil aprobado limita el universo de decisiones.",
        eyebrow_w=3.0)

    # LEFT — Diagrama conceptual de frontera eficiente
    dx = 0.85
    dy = 2.55
    dw = 6.3
    dh = 4.20
    round_rect(s, dx, dy, dw, dh, fill=WHITE, line=BORDER, radius=0.08)

    # Title del diagrama
    text(s, dx + 0.30, dy + 0.20, dw - 0.5, 0.40,
         "Frontera eficiente — concepto",
         size=13, color=NAVY_900, bold=True, font=FONT_HEAD)
    text(s, dx + 0.30, dy + 0.55, dw - 0.5, 0.30,
         "Riesgo en el eje X · Retorno esperado en el eje Y",
         size=10, color=TEXT_MUTED, italic=True, font=FONT_BODY)

    # Plot area
    px = dx + 0.80
    py = dy + 1.20
    pw = dw - 1.10
    ph = dh - 1.50

    # Axes
    hline(s, px, py + ph, px + pw, py + ph, color=TEXT_MUTED, weight=1.0)  # x-axis
    hline(s, px, py, px, py + ph, color=TEXT_MUTED, weight=1.0)            # y-axis
    text(s, px + pw - 1.0, py + ph + 0.06, 1.0, 0.25,
         "Riesgo →", size=10, color=TEXT_MUTED, align="right", italic=True, font=FONT_BODY)
    text(s, px - 0.50, py - 0.20, 2.0, 0.25,
         "↑ Retorno esperado", size=10, color=TEXT_MUTED, italic=True, font=FONT_BODY)

    # Frontier curve (a series of dots forming a concave curve)
    points = []
    for i in range(20):
        t = i / 19.0
        cx_ = px + 0.20 + t * (pw - 0.40)
        # Concave curve y = sqrt-like
        cy_ = (py + ph) - 0.25 - (ph - 0.50) * math.sqrt(t)
        points.append((cx_, cy_))
        oval(s, cx_ - 0.045, cy_ - 0.045, 0.09, fill=NAVY_600)

    # "Approved by risk budget" shaded region (translucent emerald rect)
    # Boundary: x = 0.42 * pw from left (a vertical line for max risk allowed)
    budget_x = px + 0.42 * pw
    # Soft emerald band
    rect(s, px, py + 0.10, budget_x - px, ph - 0.20,
         fill=EMERALD_LT)
    # Re-draw frontier dots on top so they remain visible
    for (cx_, cy_) in points:
        oval(s, cx_ - 0.045, cy_ - 0.045, 0.09, fill=NAVY_600)

    # Vertical line marking the risk budget limit
    hline(s, budget_x, py, budget_x, py + ph, color=EMERALD, weight=1.5)
    text(s, budget_x - 1.0, py + ph + 0.06, 2.0, 0.25,
         "presupuesto de riesgo aprobado",
         size=9, color=EMERALD, bold=True, italic=True, font=FONT_BODY,
         align="left")

    # Selected portfolio marker (a hollow ring at one of the dots within the band)
    sel_idx = 5
    scx, scy = points[sel_idx]
    oval(s, scx - 0.10, scy - 0.10, 0.20, fill=WHITE, line=EMERALD)
    oval(s, scx - 0.05, scy - 0.05, 0.10, fill=EMERALD)
    text(s, scx + 0.18, scy - 0.16, 2.5, 0.30,
         "BALANCED · seleccionada", size=9.5, color=EMERALD, bold=True, font=FONT_BODY)

    # Out-of-budget portfolio marker
    out_idx = 14
    ocx, ocy = points[out_idx]
    oval(s, ocx - 0.10, ocy - 0.10, 0.20, fill=WHITE, line=AMBER)
    oval(s, ocx - 0.05, ocy - 0.05, 0.10, fill=AMBER)
    text(s, ocx - 2.30, ocy - 0.16, 2.20, 0.30,
         "GROWTH · requiere override", size=9.5, color=AMBER, bold=True,
         align="right", font=FONT_BODY)

    # RIGHT — 4 bullets
    rx = dx + dw + 0.30
    rw = SW - rx - 0.85

    text(s, rx, 2.55, rw, 0.40,
         "TRADE-OFF Y RESTRICCIONES",
         size=11, color=NAVY_600, bold=True, charspace=180, font=FONT_BODY)
    text(s, rx, 2.95, rw, 0.50,
         "Cinco principios que el sistema respeta",
         size=18, color=NAVY_900, bold=True, font=FONT_HEAD)

    principles = [
        ("Toda cartera implica un trade-off",
         "Riesgo y retorno no se pueden optimizar por separado."),
        ("La optimización opera bajo restricciones",
         "Universo elegible, límites por activo, presupuesto de riesgo."),
        ("No toda eficiente es adecuada",
         "La frontera teórica incluye carteras fuera del perfil aprobado."),
        ("El perfil aprobado define el espacio",
         "Las variantes propuestas se comparan contra el presupuesto."),
        ("Excederlo exige justificación explícita",
         "Es la única vía para salir del espacio aprobado."),
    ]
    for i, (title, body) in enumerate(principles):
        y_ = 3.70 + i * 0.62
        oval(s, rx, y_ + 0.08, 0.12, fill=CYAN_500)
        text(s, rx + 0.28, y_, rw - 0.30, 0.30, title,
             size=12, color=NAVY_900, bold=True, font=FONT_HEAD)
        text(s, rx + 0.28, y_ + 0.27, rw - 0.30, 0.36, body,
             size=10, color=TEXT_MUTED, italic=True, font=FONT_BODY)

    footer(s, 5)


# ─────────────────────────────────────────────────────────────────────
# Slide 6 — Problema de la IA en asesoría
# ─────────────────────────────────────────────────────────────────────
def slide_6_ia_problema():
    s = add_slide(bg=BG_LIGHT)

    slide_title_block(s,
        "MARCO TEÓRICO  ·  IV",
        "La IA puede asistir, pero no debe decidir",
        sub="La responsabilidad profesional no es delegable. "
            "El sistema correcto es human-in-the-loop.",
        eyebrow_w=3.0)

    # 3-box diagram: IA → Asesor → Sistema/Auditoría
    by = 2.55
    bh = 3.05
    bw = 3.65
    bx0 = 0.85
    gap = 0.45

    boxes = [
        ("LA IA",            CYAN_500,  CYAN_LIGHT,
         "Lee, sintetiza, propone",
         "Detecta inconsistencias en el KYC. Propone un perfil preliminar. Sugiere preguntas de seguimiento.",
         "No firma. No selecciona. No decide."),
        ("EL ASESOR",        EMERALD,   EMERALD_LT,
         "Revisa, modifica, aprueba",
         "Lee la propuesta, la confronta con su criterio profesional, y la aprueba — o la modifica, o la rechaza.",
         "Es siempre el actor vinculante."),
        ("EL SISTEMA",       VIOLET,    VIOLET_LT,
         "Registra y reconstruye",
         "Guarda qué propuso la IA, qué decidió el asesor, qué cambió y por qué.",
         "La decisión queda reconstruible."),
    ]

    for i, (title, accent, accent_lt, kicker, body, foot) in enumerate(boxes):
        x = bx0 + i * (bw + gap)
        round_rect(s, x, by, bw, bh, fill=WHITE, line=BORDER, radius=0.08)
        # Top accent strip
        rect(s, x, by, bw, 0.10, fill=accent)
        # Title row with circle
        d = 0.50
        oval(s, x + 0.30, by + 0.35, d, fill=accent_lt, line=accent)
        text(s, x + 0.95, by + 0.30, bw - 1.10, 0.60, title,
             size=15, color=NAVY_900, bold=True, font=FONT_HEAD,
             charspace=120, valign="middle")
        # Kicker
        text(s, x + 0.30, by + 1.00, bw - 0.5, 0.40,
             kicker, size=12.5, color=accent, bold=True, italic=True,
             font=FONT_BODY)
        # Body
        text(s, x + 0.30, by + 1.50, bw - 0.5, 1.20,
             body, size=11, color=TEXT_BODY, font=FONT_BODY)
        # Foot
        text(s, x + 0.30, by + 2.65, bw - 0.5, 0.35,
             foot, size=11, color=NAVY_900, bold=True, italic=True,
             font=FONT_BODY)
        # Arrow to next box
        if i < len(boxes) - 1:
            arrow(s, x + bw + 0.08, by + bh / 2, x + bw + gap - 0.08,
                  by + bh / 2, color=TEXT_FAINT, weight=1.5)

    # Quote
    text(s, 0.85, 6.30, 11.6, 0.65,
         "“La IA no reemplaza el criterio profesional: lo documenta y lo potencia.”",
         size=17, color=NAVY_900, italic=True, align="center",
         font=FONT_QUOTE, bold=True)

    footer(s, 6)


# ─────────────────────────────────────────────────────────────────────
# Slide 7 — Solución propuesta
# ─────────────────────────────────────────────────────────────────────
def slide_7_solucion():
    s = add_slide(bg=BG_LIGHT)

    slide_title_block(s,
        "SOLUCIÓN PROPUESTA",
        "Risk-First Advisory: IA supervisada por asesor",
        sub="Un workflow de nueve pasos donde cada etapa tiene un responsable y "
            "queda registrada.",
        eyebrow_w=2.5)

    # 9-step flow — single row of small cards with arrows between
    steps = [
        ("Perfil del\ninversor",      NAVY_900,    "Asesor"),
        ("KYC\nestructurado",         NAVY_900,    "Asesor"),
        ("Análisis\nde IA",           CYAN_500,    "IA"),
        ("Aprobación\ndel asesor",    EMERALD,     "Asesor"),
        ("Universo\nelegible",        NAVY_700,    "Sistema"),
        ("Propuesta\nde cartera",     NAVY_700,    "Sistema"),
        ("Selección\nfinal",          EMERALD,     "Asesor"),
        ("Reporte\njustificado",      VIOLET,      "Sistema"),
        ("Auditoría\nverificada",     VIOLET,      "Compliance"),
    ]
    n = len(steps)
    row_y = 2.85
    pad_x = 0.85
    avail = SW - 2 * pad_x
    gap = 0.10
    cw = (avail - gap * (n - 1)) / n
    ch = 1.85

    # Subtle horizontal line connecting all cards
    hline(s, pad_x + 0.2, row_y + ch / 2, SW - pad_x - 0.2,
          row_y + ch / 2, color=BORDER_STR, weight=1.0)

    for i, (label, accent, owner) in enumerate(steps):
        x = pad_x + i * (cw + gap)
        round_rect(s, x, row_y, cw, ch, fill=WHITE, line=BORDER, radius=0.10)
        rect(s, x, row_y, cw, 0.08, fill=accent)
        # Number circle
        d = 0.40
        numbered_circle(s, x + cw / 2 - d / 2, row_y + 0.25, d, i + 1,
                        bg=accent, fg=WHITE, size=12)
        # Label
        text(s, x + 0.06, row_y + 0.75, cw - 0.12, 0.85,
             label, size=11, color=NAVY_900, bold=True,
             align="center", valign="middle", font=FONT_HEAD)
        # Owner
        text(s, x + 0.06, row_y + ch - 0.30, cw - 0.12, 0.25,
             owner.upper(), size=8.5, color=accent, bold=True,
             align="center", charspace=120, font=FONT_BODY)

    # Six principles row
    py_ = 5.10
    principles = [
        ("LA IA",          "propone"),
        ("EL ASESOR",      "aprueba, modifica o rechaza"),
        ("EL SISTEMA",     "controla risk budget y suitability"),
        ("LA CARTERA",     "muestra instrumentos y pesos"),
        ("EL REPORTE",     "queda justificado y revisable"),
        ("EL PROCESO",     "queda auditado de extremo a extremo"),
    ]
    pcols = 3
    pcw = (SW - 2 * pad_x - (pcols - 1) * 0.30) / pcols
    pch = 0.75
    pgap_y = 0.18
    for i, (k, v) in enumerate(principles):
        col = i % pcols
        row = i // pcols
        x = pad_x + col * (pcw + 0.30)
        y = py_ + row * (pch + pgap_y)
        oval(s, x, y + 0.30, 0.14, fill=NAVY_600)
        text(s, x + 0.25, y + 0.10, pcw - 0.30, 0.30,
             k, size=11, color=NAVY_900, bold=True, charspace=140, font=FONT_BODY)
        text(s, x + 0.25, y + 0.38, pcw - 0.30, 0.30,
             v, size=12, color=TEXT_BODY, italic=True, font=FONT_BODY)

    footer(s, 7)


# ─────────────────────────────────────────────────────────────────────
# Slide 8 — Demo funcional actual
# ─────────────────────────────────────────────────────────────────────
def slide_8_demo():
    s = add_slide(bg=BG_LIGHT)

    slide_title_block(s,
        "PRODUCTO  ·  DEMO",
        "Qué permite mostrar hoy la demo",
        sub="Un prototipo funcional que recorre el workflow completo "
            "en una máquina local, en español, paso a paso.",
        eyebrow_w=2.7)

    items = [
        ("Demo en español",
         "Hero, formulario, instrucciones y mensajes localizados para el asesor."),
        ("Carga de perfil inversor",
         "Datos básicos, preguntas abiertas y preferencias en un solo formulario."),
        ("Flujo paso a paso",
         "Ocho pasos didácticos con explicación, acción y resultado por paso."),
        ("Propuesta de cartera",
         "Variantes DEFENSIVE / BALANCED / GROWTH con instrumentos y pesos visibles."),
        ("Reporte Markdown",
         "Vista previa del documento que el asesor revisaría antes de compartirlo."),
        ("Auditoría del caso",
         "Verificación de la integridad del registro de decisiones del caso."),
        ("Logs de IA con redacción",
         "Cada llamada al modelo guardada con el texto del cliente sanitizado."),
        ("Bootstrap local",
         "Un solo comando deja el entorno listo para correr la demo desde cero."),
    ]
    cols = 4
    pad_x = 0.85
    cy0 = 2.55
    cw = (SW - 2 * pad_x - (cols - 1) * 0.20) / cols
    ch = 1.95
    gap_y = 0.20

    for i, (title, body) in enumerate(items):
        col = i % cols
        row = i // cols
        x = pad_x + col * (cw + 0.20)
        y = cy0 + row * (ch + gap_y)
        round_rect(s, x, y, cw, ch, fill=WHITE, line=BORDER, radius=0.08)
        # Small check icon
        oval(s, x + 0.25, y + 0.30, 0.30, fill=EMERALD_LT, line=EMERALD)
        text(s, x + 0.25, y + 0.30, 0.30, 0.30, "✓",
             size=12, color=EMERALD, bold=True, align="center", valign="middle",
             font=FONT_BODY)
        text(s, x + 0.65, y + 0.27, cw - 0.80, 0.45,
             title, size=12.5, color=NAVY_900, bold=True, font=FONT_HEAD)
        text(s, x + 0.25, y + 0.90, cw - 0.45, ch - 1.05,
             body, size=10.5, color=TEXT_BODY, font=FONT_BODY)

    # Disclaimer footer
    round_rect(s, 0.85, 6.80, 11.6, 0.42, fill=AMBER_LT, line=AMBER,
               radius=0.20, line_w=0.5)
    text(s, 0.85, 6.80, 11.6, 0.42,
         "El universo de instrumentos es acotado / demo. "
         "No es market data productivo.",
         size=11.5, color=AMBER, italic=True, align="center", valign="middle",
         font=FONT_BODY, bold=True)


# ─────────────────────────────────────────────────────────────────────
# Slide 9 — Propuesta de cartera transparente
# ─────────────────────────────────────────────────────────────────────
def slide_9_propuesta_transparente():
    s = add_slide(bg=BG_LIGHT)

    slide_title_block(s,
        "PROPUESTA DE CARTERA",
        "La cartera no es una etiqueta: muestra instrumentos, pesos y motivos",
        sub="Una propuesta útil para el asesor — y para el cliente — debe "
            "responder qué se compra, cuánto y por qué.",
        eyebrow_w=2.7)

    # LEFT — variant card with mockup table (BALANCED)
    lx = 0.85
    ly = 2.50
    lw = 7.7
    lh = 4.55
    round_rect(s, lx, ly, lw, lh, fill=WHITE, line=BORDER, radius=0.08)
    rect(s, lx, ly, lw, 0.10, fill=EMERALD)

    text(s, lx + 0.30, ly + 0.25, lw - 0.60, 0.45,
         "BALANCED  —  dentro del presupuesto de riesgo",
         size=14, color=NAVY_900, bold=True, font=FONT_HEAD)
    text(s, lx + 0.30, ly + 0.72, lw - 0.60, 0.30,
         "Retorno esperado 6,8 %   ·   Volatilidad 9,4 %   ·   6 instrumentos",
         size=11, color=TEXT_MUTED, italic=True, font=FONT_BODY)

    # Table headers
    hy = ly + 1.15
    rh = 0.42
    col_x = [lx + 0.30, lx + 1.85, lx + 3.30, lx + 4.55, lx + 5.20, lx + 5.95]
    col_w = [1.55, 1.45, 1.25, 0.65, 0.75, 1.55]
    headers = ["INSTRUMENTO", "TIPO", "MONEDA", "PESO", "BARRA", "MOTIVO"]
    rect(s, lx + 0.20, hy, lw - 0.40, rh, fill=BG_SOFT)
    for i, hd in enumerate(headers):
        if hd == "BARRA":
            continue
        align = "right" if hd == "PESO" else "left"
        text(s, col_x[i], hy, col_w[i], rh, hd,
             size=9, color=TEXT_MUTED, bold=True, valign="middle",
             align=align, charspace=120, font=FONT_BODY)

    rows = [
        ("GD30", "Bono soberano",     "USD", 0.30, "Renta fija · hard dollar"),
        ("AL30", "Bono soberano",     "USD", 0.22, "Renta fija · hard dollar"),
        ("YPFD", "Acción",            "ARS", 0.15, "Equity · energía"),
        ("SHV",  "ETF money market",  "USD", 0.13, "Liquidez · USD"),
        ("PAMP", "CEDEAR",            "USD", 0.10, "Equity · diversificación"),
        ("MEP",  "Money market USD",  "USD", 0.10, "Liquidez · USD"),
    ]
    max_w = max(w for _, _, _, w, _ in rows)
    for i, (ticker, kind, cur, w, motivo) in enumerate(rows):
        ry = hy + rh + 0.05 + i * (rh + 0.05)
        rect(s, lx + 0.20, ry, lw - 0.40, rh, fill=WHITE)
        if i < len(rows) - 1:
            hline(s, lx + 0.30, ry + rh, lx + lw - 0.30,
                  ry + rh, color=BORDER, weight=0.4)
        text(s, col_x[0], ry, col_w[0], rh, ticker,
             size=12, color=NAVY_900, bold=True, font=FONT_MONO, valign="middle")
        text(s, col_x[1], ry, col_w[1], rh, kind,
             size=11, color=TEXT_BODY, valign="middle", font=FONT_BODY)
        text(s, col_x[2], ry, col_w[2], rh, cur,
             size=11, color=TEXT_BODY, valign="middle", font=FONT_BODY)
        text(s, col_x[3], ry, col_w[3], rh, f"{int(w*100)} %",
             size=11.5, color=NAVY_900, bold=True, valign="middle",
             align="right", font=FONT_BODY)
        bar_x = col_x[4]
        bar_w = col_w[4]
        bar_y = ry + rh / 2 - 0.07
        rect(s, bar_x, bar_y, bar_w, 0.14, fill=BG_SOFT)
        fill_w = bar_w * (w / max_w)
        rect(s, bar_x, bar_y, fill_w, 0.14, fill=NAVY_600)
        text(s, col_x[5], ry, col_w[5], rh, motivo,
             size=10.5, color=TEXT_MUTED, valign="middle", italic=True,
             font=FONT_BODY)

    # RIGHT — 5 preguntas que la propuesta responde
    rx = lx + lw + 0.30
    rw = SW - rx - 0.85

    text(s, rx, 2.55, rw, 0.40,
         "UNA PROPUESTA ÚTIL RESPONDE",
         size=11, color=NAVY_600, bold=True, charspace=180, font=FONT_BODY)
    text(s, rx, 2.95, rw, 0.50,
         "Cinco preguntas concretas",
         size=18, color=NAVY_900, bold=True, font=FONT_HEAD)

    qs = [
        ("Qué compro",        "Cada instrumento explicitado por ticker y nombre."),
        ("Cuánto compro",     "Peso porcentual y barra visual para comparar."),
        ("Por qué entra",     "Motivo: clase de activo, moneda, función en la cartera."),
        ("Qué riesgo implica", "Tipo de instrumento + perfil al que pertenece."),
        ("Si excede el presupuesto", "Marcador claro si requiere override del asesor."),
    ]
    for i, (q, a) in enumerate(qs):
        y_ = 3.70 + i * 0.62
        oval(s, rx, y_ + 0.10, 0.14, fill=CYAN_500)
        text(s, rx + 0.30, y_, rw - 0.30, 0.30, q,
             size=12, color=NAVY_900, bold=True, font=FONT_HEAD)
        text(s, rx + 0.30, y_ + 0.27, rw - 0.30, 0.36, a,
             size=10, color=TEXT_MUTED, italic=True, font=FONT_BODY)

    footer(s, 9)


# ─────────────────────────────────────────────────────────────────────
# Slide 10 — Auditoría, trazabilidad y responsabilidad
# ─────────────────────────────────────────────────────────────────────
def slide_10_trazabilidad():
    s = add_slide(bg=BG_LIGHT)

    slide_title_block(s,
        "TRAZABILIDAD",
        "De recomendación opaca a decisión reconstruible",
        sub="Toda recomendación profesional debería poder reconstruirse "
            "desde la data original hasta el reporte final.",
        eyebrow_w=2.0)

    # Layout: vertical timeline of "qué queda registrado"
    items = [
        ("01", "Qué datos entraron",
         "El KYC estructurado y las respuestas en texto libre del cliente."),
        ("02", "Qué perfil propuso la IA",
         "Perfil preliminar sugerido + contradicciones detectadas."),
        ("03", "Qué aprobó el asesor",
         "Perfil final aprobado, decisión (approve/modify/reject) y rationale."),
        ("04", "Qué universo se usó",
         "El conjunto de instrumentos elegibles tras aplicar las preferencias."),
        ("05", "Qué cartera se propuso",
         "Las tres variantes generadas, con instrumentos, pesos y métricas."),
        ("06", "Qué se seleccionó",
         "La variante final y el override firmado por el asesor (si aplica)."),
        ("07", "Qué reporte se generó",
         "Markdown determinístico con composición y comparación de variantes."),
    ]
    n = len(items)
    pad_x = 0.85
    cy0 = 2.55
    avail_h = 7.10 - cy0 - 0.50  # leave room for footer + small note
    h_per = avail_h / n

    tx = pad_x + 0.25
    timeline_x = tx + 0.18 / 2

    # Vertical timeline
    hline(s, timeline_x, cy0 + 0.18, timeline_x,
          cy0 + (n - 1) * h_per + 0.30, color=BORDER_STR, weight=2.0)

    for i, (n_, title, body) in enumerate(items):
        y = cy0 + i * h_per
        # Numbered dot
        oval(s, tx, y + 0.10, 0.28, fill=NAVY_600)
        text(s, tx, y + 0.10, 0.28, 0.28, n_,
             size=10, color=WHITE, bold=True, align="center", valign="middle",
             font=FONT_BODY)
        cx = tx + 0.55
        cw = SW - cx - 0.85
        text(s, cx, y + 0.07, cw, 0.32, title,
             size=13, color=NAVY_900, bold=True, font=FONT_HEAD)
        text(s, cx, y + 0.36, cw, 0.30, body,
             size=10.5, color=TEXT_BODY, font=FONT_BODY)

    # Small implementation note at the bottom (tech detail kept as backup)
    round_rect(s, 0.85, 7.00, 11.6, 0.32, fill=BG_SOFT, line=BORDER,
               radius=0.30, line_w=0.5)
    text(s, 0.85, 7.00, 11.6, 0.32,
         "Implementación: cadena hash SHA-256 por caso y logs de IA "
         "con PII redactada por el backend.",
         size=10, color=TEXT_MUTED, italic=True, align="center", valign="middle",
         font=FONT_BODY)


# ─────────────────────────────────────────────────────────────────────
# Slide 11 — Estado actual y limitaciones
# ─────────────────────────────────────────────────────────────────────
def slide_11_estado_y_limites():
    s = add_slide(bg=BG_LIGHT)

    slide_title_block(s,
        "ESTADO ACTUAL",
        "Prototipo funcional, no producto final",
        sub="Honestidad sobre el alcance: lo que ya funciona y lo que "
            "queda pendiente para pasar a piloto.",
        eyebrow_w=1.9)

    # Two columns: works / limitations
    col_w = 5.7
    col_h = 4.10
    cy = 2.55
    left_x = 0.85
    right_x = left_x + col_w + 0.30

    # LEFT — Ya funciona
    card(s, left_x, cy, col_w, col_h, fill=WHITE, accent=EMERALD, accent_w=0.06)
    text(s, left_x + 0.35, cy + 0.25, col_w - 0.5, 0.35,
         "YA FUNCIONA", size=11, color=EMERALD, bold=True,
         charspace=180, font=FONT_BODY)
    text(s, left_x + 0.35, cy + 0.65, col_w - 0.5, 0.50,
         "Lo que la demo puede mostrar hoy",
         size=16, color=NAVY_900, bold=True, font=FONT_HEAD)
    works = [
        "Workflow case-scoped completo de extremo a extremo",
        "Demo local plug-and-play",
        "Flujo investor-facing en español, paso a paso",
        "Propuestas de cartera con instrumentos y pesos",
        "Reporte Markdown justificado",
        "Auditoría del caso (cadena de decisiones íntegra)",
        "Tests de regresión que cubren el flujo completo",
    ]
    for i, w in enumerate(works):
        y_ = cy + 1.40 + i * 0.36
        oval(s, left_x + 0.45, y_ + 0.13, 0.10, fill=EMERALD)
        text(s, left_x + 0.68, y_, col_w - 0.85, 0.32,
             w, size=11.5, color=TEXT_BODY, font=FONT_BODY)

    # RIGHT — Limitaciones
    card(s, right_x, cy, col_w, col_h, fill=WHITE, accent=AMBER, accent_w=0.06)
    text(s, right_x + 0.35, cy + 0.25, col_w - 0.5, 0.35,
         "LIMITACIONES", size=11, color=AMBER, bold=True,
         charspace=180, font=FONT_BODY)
    text(s, right_x + 0.35, cy + 0.65, col_w - 0.5, 0.50,
         "Lo que aún no — y por qué importa",
         size=16, color=NAVY_900, bold=True, font=FONT_HEAD)
    limits = [
        "No production-ready — demo local, no piloto",
        "Universo de instrumentos acotado / demo",
        "Sin market data provider real",
        "Autenticación de desarrollo, no productiva",
        "Reporte solo Markdown, sin PDF / branding",
        "Sin firm-level access control completo",
        "No usar con datos reales sensibles",
    ]
    for i, limit_text in enumerate(limits):
        y_ = cy + 1.40 + i * 0.36
        oval(s, right_x + 0.45, y_ + 0.13, 0.10, fill=AMBER)
        text(s, right_x + 0.68, y_, col_w - 0.85, 0.32,
             limit_text, size=11.5, color=TEXT_BODY, font=FONT_BODY)

    # Small backup note: tech credibility
    round_rect(s, 0.85, 6.95, 11.6, 0.35, fill=BG_SOFT, line=BORDER,
               radius=0.25, line_w=0.5)
    text(s, 0.85, 6.95, 11.6, 0.35,
         "Respaldo técnico (no protagonista): backend case-scoped, "
         "auditoría implementada, ~3 087 tests automatizados verdes.",
         size=10, color=TEXT_MUTED, italic=True, align="center", valign="middle",
         font=FONT_BODY)


# ─────────────────────────────────────────────────────────────────────
# Slide 12 — Roadmap y cierre
# ─────────────────────────────────────────────────────────────────────
def slide_12_cierre():
    s = add_slide(bg=NAVY_DARK)

    # Decoración
    for x in [0.55, 0.75]:
        hline(s, x, 0.0, x, SH, color=NAVY_700, weight=0.5)

    # Brand mark
    oval(s, 1.10, 0.78, 0.18, fill=CYAN_500)
    text(s, 1.40, 0.66, 6.0, 0.5,
         "RISK-FIRST  ADVISORY", size=11, color=SKY_LIGHT, bold=True,
         charspace=240, font=FONT_BODY)

    # Title
    text(s, 1.10, 1.40, 11.5, 0.60,
         "Próximo paso: de prototipo académico a piloto controlado",
         size=26, color=WHITE, bold=True, font=FONT_HEAD)

    # LEFT — Roadmap items
    rx = 1.10
    ry = 2.35
    rw = 6.5
    text(s, rx, ry, rw, 0.40,
         "ROADMAP  ·  FASE 4",
         size=11, color=CYAN_LIGHT, bold=True, charspace=180, font=FONT_BODY)

    roadmap = [
        "Market data freshness",
        "Carga manual del universo de instrumentos",
        "PDF / export con branding",
        "Access control por firma",
        "Autenticación productiva",
        "Backup y restore",
        "Piloto con asesor real",
    ]
    for i, item in enumerate(roadmap):
        y_ = ry + 0.50 + i * 0.42
        oval(s, rx, y_ + 0.12, 0.16, fill=CYAN_500)
        text(s, rx + 0.30, y_, rw - 0.30, 0.34, item,
             size=14, color=LIGHT_INK, font=FONT_BODY)

    # RIGHT — Cierre / pitch
    cx = 7.85
    cw = SW - cx - 0.85
    round_rect(s, cx, ry, cw, 4.20, fill=NAVY_800, line=NAVY_700,
               radius=0.08)
    text(s, cx + 0.35, ry + 0.25, cw - 0.5, 0.35,
         "MENSAJE FINAL", size=11, color=CYAN_LIGHT, bold=True,
         charspace=180, font=FONT_BODY)

    text_lines(s, cx + 0.35, ry + 0.70, cw - 0.5, 2.40, [
        [("Risk-First Advisory busca convertir la asesoría asistida por IA en un proceso ",
          {"size": 14, "color": LIGHT_INK, "font": FONT_BODY})],
        [("controlada, justificable y auditable",
          {"size": 14, "color": WHITE, "font": FONT_BODY, "bold": True})],
        [(", donde la tecnología potencia al asesor",
          {"size": 14, "color": LIGHT_INK, "font": FONT_BODY})],
        [(" sin reemplazarlo.",
          {"size": 14, "color": CYAN_500, "font": FONT_BODY, "bold": True})],
    ], line_spacing=1.35, space_after=2)

    # Three final taglines
    hline(s, cx + 0.35, ry + 2.85, cx + cw - 0.35, ry + 2.85,
          color=NAVY_700, weight=0.75)
    text_lines(s, cx + 0.35, ry + 2.95, cw - 0.5, 1.20, [
        [("“La IA propone. El asesor decide.”",
          {"size": 13, "color": CYAN_LIGHT, "italic": True, "font": FONT_QUOTE})],
        [("“Risk-first, not return-chasing.”",
          {"size": 13, "color": CYAN_LIGHT, "italic": True, "font": FONT_QUOTE})],
        [("“Toda decisión debe poder justificarse.”",
          {"size": 13, "color": CYAN_LIGHT, "italic": True, "font": FONT_QUOTE})],
    ], line_spacing=1.40, space_after=2)

    # Bottom row
    by = 6.55
    bh = 0.75
    bw = (SW - 2.20) / 2
    items = [
        ("ESTADO",         "Prototipo funcional · demo local"),
        ("PRÓXIMO HITO",   "Fase 4 — piloto controlado con asesor real"),
    ]
    for i, (label, val) in enumerate(items):
        x = 1.10 + i * (bw + 0.20)
        round_rect(s, x, by, bw, bh, fill=NAVY_800, line=NAVY_700,
                   radius=0.10)
        text(s, x + 0.25, by + 0.12, bw - 0.5, 0.30,
             label, size=10, color=CYAN_LIGHT, bold=True, charspace=160,
             font=FONT_BODY)
        text(s, x + 0.25, by + 0.40, bw - 0.5, 0.30,
             val, size=12, color=WHITE, bold=True, font=FONT_BODY)


# ─────────────────────────────────────────────────────────────────────
# Build
# ─────────────────────────────────────────────────────────────────────
def main():
    slide_1_portada()
    slide_2_problema()
    slide_3_suitability()
    slide_4_risk_first()
    slide_5_carteras_restricciones()
    slide_6_ia_problema()
    slide_7_solucion()
    slide_8_demo()
    slide_9_propuesta_transparente()
    slide_10_trazabilidad()
    slide_11_estado_y_limites()
    slide_12_cierre()

    out = "Risk-First-Advisory-Academic-Deck.pptx"
    prs.save(out)
    print(f"Wrote {out} ({len(prs.slides)} slides, "
          f"{prs.slide_width.inches:.3f}x{prs.slide_height.inches:.3f} in)")


if __name__ == "__main__":
    main()
