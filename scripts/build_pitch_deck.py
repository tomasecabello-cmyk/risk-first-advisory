#!/usr/bin/env python3
"""
Build Risk-First Advisory pitch deck (12 slides, 16:9 widescreen).

Aesthetic: institutional fintech / wealth-management. Palette matches the
actual product UI (navy / cyan / emerald / amber). Designed for a
non-technical finance audience (profesor, asesor, mentor).

Run: python scripts/build_pitch_deck.py
Output: Risk-First-Advisory-Pitch-Deck.pptx
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ─────────────────────────────────────────────────────────────────────
# Palette  (mirrors the Risk-First Advisory product design tokens)
# ─────────────────────────────────────────────────────────────────────
NAVY_DARK  = RGBColor(0x0A, 0x1A, 0x36)   # hero bg deepest
NAVY_900   = RGBColor(0x0B, 0x1E, 0x3F)   # brand
NAVY_800   = RGBColor(0x12, 0x2A, 0x4F)
NAVY_700   = RGBColor(0x1E, 0x3A, 0x6B)
NAVY_600   = RGBColor(0x1D, 0x4E, 0xD8)   # primary action
NAVY_500   = RGBColor(0x25, 0x63, 0xEB)
CYAN_500   = RGBColor(0x0E, 0xA5, 0xE9)   # accent (AI / data)
CYAN_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
SKY_LIGHT  = RGBColor(0xBF, 0xDB, 0xFE)
EMERALD    = RGBColor(0x05, 0x96, 0x69)
EMERALD_LT = RGBColor(0xD1, 0xFA, 0xE5)
AMBER      = RGBColor(0xD9, 0x77, 0x06)
AMBER_LT   = RGBColor(0xFE, 0xF3, 0xC7)
ROSE       = RGBColor(0xDC, 0x26, 0x26)
VIOLET     = RGBColor(0x6D, 0x28, 0xD9)
VIOLET_LT  = RGBColor(0xED, 0xE9, 0xFE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF8, 0xFA, 0xFC)
BG_LIGHT   = RGBColor(0xF4, 0xF6, 0xFB)
BG_SOFT    = RGBColor(0xEE, 0xF1, 0xF7)
BORDER     = RGBColor(0xE5, 0xE7, 0xEB)
BORDER_STR = RGBColor(0xCB, 0xD5, 0xE1)
TEXT_DARK  = RGBColor(0x0F, 0x17, 0x2A)
TEXT_BODY  = RGBColor(0x33, 0x41, 0x55)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
TEXT_FAINT = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_INK  = RGBColor(0xE2, 0xE8, 0xF0)   # text on dark bg
SOFT_INK   = RGBColor(0xCB, 0xD5, 0xE1)

FONT_HEAD  = "Calibri"
FONT_BODY  = "Calibri"
FONT_MONO  = "Consolas"


# ─────────────────────────────────────────────────────────────────────
# Presentation setup (16:9 widescreen)
# ─────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
prs.core_properties.title  = "Risk-First Advisory — Pitch deck"
prs.core_properties.author = "Risk-First Advisory"
BLANK = prs.slide_layouts[6]

SW = 13.333  # slide width in inches
SH = 7.5     # slide height in inches


# ─────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────

def add_slide(bg=BG_LIGHT):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                           prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def rect(slide, x, y, w, h, *, fill=WHITE, line=None, line_w=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is not None:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def round_rect(slide, x, y, w, h, *, fill=WHITE, line=None, line_w=0.5, radius=0.10):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    s.adjustments[0] = radius
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is not None:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def oval(slide, x, y, d, *, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(x), Inches(y), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is not None:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def hline(slide, x1, y1, x2, y2, *, color=BORDER, weight=0.5):
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1),
                                    Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def _set_char_spacing(run, hundredths_pt: int):
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(hundredths_pt)))


def text(slide, x, y, w, h, body, *, size=14, color=TEXT_BODY, bold=False,
         italic=False, align="left", valign="top", font=FONT_BODY, charspace=None):
    """body: str | list of (str | (str, dict))."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]

    def write_run(p_, seg, default_opts):
        if isinstance(seg, str):
            seg_text, seg_opts = seg, {}
        else:
            seg_text, seg_opts = seg
        r_ = p_.add_run()
        r_.text = seg_text
        f = r_.font
        f.name   = seg_opts.get("font", default_opts["font"])
        f.size   = Pt(seg_opts.get("size", default_opts["size"]))
        f.bold   = seg_opts.get("bold", default_opts["bold"])
        f.italic = seg_opts.get("italic", default_opts["italic"])
        f.color.rgb = seg_opts.get("color", default_opts["color"])
        cs = seg_opts.get("charspace", default_opts["charspace"])
        if cs:
            _set_char_spacing(r_, cs)

    defaults = {"size": size, "color": color, "bold": bold, "italic": italic,
                "font": font, "charspace": charspace}

    if isinstance(body, list):
        for seg in body:
            write_run(p, seg, defaults)
    else:
        write_run(p, body, defaults)
    return tb


def text_lines(slide, x, y, w, h, lines, *, size=14, color=TEXT_BODY, bold=False,
               align="left", valign="top", font=FONT_BODY, line_spacing=1.18,
               space_after=4):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM

    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}[align]
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        segs = ln if isinstance(ln, list) else [ln]
        for seg in segs:
            if isinstance(seg, str):
                seg_text, seg_opts = seg, {}
            else:
                seg_text, seg_opts = seg
            r_ = p.add_run()
            r_.text = seg_text
            f = r_.font
            f.name   = seg_opts.get("font", font)
            f.size   = Pt(seg_opts.get("size", size))
            f.bold   = seg_opts.get("bold", bold)
            f.italic = seg_opts.get("italic", False)
            f.color.rgb = seg_opts.get("color", color)
    return tb


def eyebrow_chip(slide, x, y, label, *, fg=WHITE, bg=NAVY_900, w=2.3, h=0.34):
    """Small rectangular eyebrow chip (anti-AI-slop: no underline)."""
    round_rect(slide, x, y, w, h, fill=bg, radius=0.45)
    text(slide, x, y, w, h, label, size=10.5, color=fg, bold=True,
         align="center", valign="middle", charspace=160)


def card(slide, x, y, w, h, *, fill=WHITE, accent=None, accent_w=0.07):
    """White rounded card with optional left-edge accent stripe."""
    round_rect(slide, x, y, w, h, fill=fill, line=BORDER, radius=0.06)
    if accent is not None:
        # Use a RECTANGLE strip overlaid on the LEFT edge — the card itself
        # has very small radius so the strip lines up cleanly.
        rect(slide, x, y, accent_w, h, fill=accent)


def numbered_circle(slide, x, y, d, num, *, bg=NAVY_900, fg=WHITE, size=14):
    oval(slide, x, y, d, fill=bg)
    text(slide, x, y, d, d, str(num), size=size, color=fg, bold=True,
         align="center", valign="middle")


def footer(slide, num, total=12, *, ink=TEXT_FAINT):
    text(slide, SW - 1.3, SH - 0.4, 1.0, 0.3,
         f"{num} / {total}", size=9, color=ink, align="right")
    text(slide, 0.55, SH - 0.4, 4.0, 0.3,
         "Risk-First Advisory · Pitch deck", size=9, color=ink, align="left")


# ─────────────────────────────────────────────────────────────────────
# Slide 1 — Portada (dark)
# ─────────────────────────────────────────────────────────────────────
def slide_1_portada():
    s = add_slide(bg=NAVY_DARK)

    # Subtle grid-like decoration: a few thin cyan vertical lines on the right side
    for i, x in enumerate([10.6, 11.4, 12.2]):
        hline(s, x, 0.0, x, SH, color=NAVY_700, weight=0.5)

    # Accent dot under brand mark (top-left)
    oval(s, 0.85, 0.78, 0.20, fill=CYAN_500)
    text(s, 1.18, 0.66, 6.0, 0.5,
         "RISK-FIRST  ADVISORY",
         size=12, color=SKY_LIGHT, bold=True, charspace=240)

    # Eyebrow
    eyebrow_chip(s, 0.85, 1.7, "WEALTH-ADVISORY BACKEND  ·  DEMO LOCAL",
                 fg=CYAN_LIGHT, bg=NAVY_800, w=4.8)

    # Hero title — 3 lines
    text_lines(s, 0.85, 2.25, 11.6, 2.6,
        [
            [("Risk-First Advisory", {"size": 60, "bold": True, "color": WHITE})],
            [("IA supervisada por asesores para perfilar clientes,",
              {"size": 26, "bold": False, "color": LIGHT_INK})],
            [("construir carteras justificables y auditar decisiones.",
              {"size": 26, "bold": False, "color": LIGHT_INK})],
        ],
        line_spacing=1.1, space_after=4)

    # English tagline (small, under the subtitle)
    text(s, 0.85, 5.45, 11.6, 0.4,
         "“AI proposes. Advisor decides. Every step audited.”",
         size=16, color=CYAN_LIGHT, italic=True)

    # Three quick-facts row
    facts = [
        ("15",          "Pasos del workflow"),
        ("SHA-256",     "Cadena de auditoría por caso"),
        ("3 087",       "Tests verdes"),
    ]
    fx = 0.85
    fw = 3.5
    fy = 6.05
    fh = 0.95
    for i, (num, label) in enumerate(facts):
        x = fx + i * (fw + 0.25)
        round_rect(s, x, fy, fw, fh, fill=NAVY_800, line=NAVY_700, radius=0.08)
        text(s, x + 0.20, fy + 0.08, fw - 0.4, 0.50, num,
             size=24, color=WHITE, bold=True)
        text(s, x + 0.20, fy + 0.55, fw - 0.4, 0.30, label,
             size=10.5, color=SOFT_INK, charspace=80)

    # Bottom-right disclaimer corner
    text(s, SW - 5.4, SH - 0.55, 4.85, 0.35,
         "Prototipo funcional · demo local · no production-ready",
         size=10, color=TEXT_FAINT, italic=True, align="right")


# ─────────────────────────────────────────────────────────────────────
# Slide 2 — Problema
# ─────────────────────────────────────────────────────────────────────
def slide_2_problema():
    s = add_slide(bg=BG_LIGHT)

    eyebrow_chip(s, 0.85, 0.6, "EL PROBLEMA HOY", fg=WHITE, bg=NAVY_900, w=2.1)
    text(s, 0.85, 1.05, 11.6, 0.9,
         "La asesoría asistida por IA tiene cinco grietas",
         size=32, color=NAVY_900, bold=True)
    text(s, 0.85, 1.85, 11.6, 0.55,
         "Cinco puntos críticos que aparecen en casi todos los procesos actuales — y que cualquier mentor o regulador puede señalar.",
         size=14, color=TEXT_MUTED)

    pains = [
        ("KYC fragmentado",
         "El perfil del cliente se arma en varios sistemas (forms, planillas, mails). Difícil de versionar y auditar."),
        ("Perfilamiento subjetivo",
         "El criterio del asesor varía. Sin reglas explícitas, dos asesores armarían perfiles distintos para el mismo cliente."),
        ("Recomendaciones poco trazables",
         "El cliente recibe una cartera pero no queda claro qué KYC la justifica ni qué supuestos usó el sistema."),
        ("Overrides sin firma formal",
         "Cuando se excede el presupuesto de riesgo aprobado, el override se discute en un email — no como decisión firmada."),
        ("Sin auditoría IA ↔ asesor ↔ cliente",
         "No hay un registro inmutable de qué propuso la IA, qué decidió el asesor y qué vio el cliente al final."),
    ]

    # 2 columns layout: 3 left, 2 right
    col_w = 5.7
    row_h = 1.4
    gap_x = 0.30
    gap_y = 0.15
    x0 = 0.85
    y0 = 2.55

    for i, (title, body) in enumerate(pains):
        col = 0 if i < 3 else 1
        row = i if col == 0 else i - 3
        x = x0 + col * (col_w + gap_x)
        y = y0 + row * (row_h + gap_y)
        card(s, x, y, col_w, row_h, fill=WHITE, accent=NAVY_600, accent_w=0.08)

        # Number badge
        numbered_circle(s, x + 0.30, y + 0.30, 0.45, i + 1,
                        bg=NAVY_900, fg=WHITE, size=14)
        # Title + body
        text(s, x + 0.95, y + 0.27, col_w - 1.2, 0.45, title,
             size=15, color=NAVY_900, bold=True)
        text(s, x + 0.95, y + 0.72, col_w - 1.2, 0.65, body,
             size=11.5, color=TEXT_BODY)

    footer(s, 2)


# ─────────────────────────────────────────────────────────────────────
# Slide 3 — Idea central
# ─────────────────────────────────────────────────────────────────────
def slide_3_idea_central():
    s = add_slide(bg=BG_LIGHT)

    eyebrow_chip(s, 0.85, 0.6, "IDEA CENTRAL", fg=WHITE, bg=NAVY_900, w=1.7)
    text(s, 0.85, 1.05, 11.6, 0.9,
         "Una división explícita de roles, auditada",
         size=32, color=NAVY_900, bold=True)

    # Big quote area
    quote_y = 2.05
    text_lines(s, 0.85, quote_y, 11.6, 1.2,
        [
            [("“La IA propone.  ", {"size": 44, "bold": True, "color": NAVY_900}),
             ("El asesor decide.”", {"size": 44, "bold": True, "color": CYAN_500})],
        ],
        line_spacing=1.0)
    text(s, 0.85, quote_y + 1.10, 11.6, 0.45,
         "“AI proposes. Advisor decides.”",
         size=15, color=TEXT_MUTED, italic=True)

    # Four roles as small cards (matches actual product story-strip)
    roles = [
        ("LA IA",         "Interpreta el KYC y propone perfil + contradicciones.",
                          "Nunca decide.",                         CYAN_500),
        ("EL ASESOR",     "Aprueba, modifica o rechaza el perfil. Firma decisiones.",
                          "Siempre vinculante.",                   EMERALD),
        ("EL SISTEMA",    "Controla suitability, presupuesto de riesgo, ESG, data quality.",
                          "Bloquea excepciones.",                  AMBER),
        ("LA AUDITORÍA",  "Hash chain SHA-256 por caso, AI logs con PII redactada.",
                          "Trazabilidad total.",                   VIOLET),
    ]
    cy = 3.95
    ch = 3.05
    cw = 2.85
    cx0 = 0.85
    gap = 0.18

    for i, (eyebrow, body, kicker, accent) in enumerate(roles):
        x = cx0 + i * (cw + gap)
        card(s, x, cy, cw, ch, fill=WHITE, accent=accent, accent_w=0.08)
        text(s, x + 0.35, cy + 0.30, cw - 0.55, 0.40, eyebrow,
             size=11, color=accent, bold=True, charspace=200)
        text(s, x + 0.35, cy + 0.85, cw - 0.55, 1.4, body,
             size=13, color=TEXT_BODY)
        text(s, x + 0.35, cy + 2.40, cw - 0.55, 0.45, kicker,
             size=13, color=NAVY_900, bold=True, italic=True)

    footer(s, 3)


# ─────────────────────────────────────────────────────────────────────
# Slide 4 — Flujo del producto (9-step timeline)
# ─────────────────────────────────────────────────────────────────────
def slide_4_flujo():
    s = add_slide(bg=BG_LIGHT)

    eyebrow_chip(s, 0.85, 0.6, "FLUJO DEL PRODUCTO", fg=WHITE, bg=NAVY_900, w=2.4)
    text(s, 0.85, 1.05, 11.6, 0.9,
         "Nueve pasos auditables, de la entrevista al reporte",
         size=30, color=NAVY_900, bold=True)
    text(s, 0.85, 1.85, 11.6, 0.50,
         "Cada paso queda registrado en el caso del cliente. El asesor avanza paso por paso; el sistema impide saltos de control.",
         size=14, color=TEXT_MUTED)

    steps = [
        ("Perfil del\ninversor",         NAVY_900),
        ("KYC\nestructurado",            NAVY_900),
        ("Análisis\nde IA",              CYAN_500),
        ("Aprobación\ndel asesor",       EMERALD),
        ("Universo\nelegible",           NAVY_700),
        ("Propuesta\nde cartera",        NAVY_700),
        ("Selección\nfinal",             EMERALD),
        ("Reporte\nrevisable",           VIOLET),
        ("Auditoría\nverificada",        VIOLET),
    ]

    # Single row of 9 mini-cards
    n = len(steps)
    row_y = 2.95
    pad_x = 0.85
    avail = SW - 2 * pad_x
    gap = 0.10
    cw = (avail - gap * (n - 1)) / n
    ch = 2.20

    # Background timeline line (subtle)
    hline(s, pad_x + 0.2, row_y + ch / 2, SW - pad_x - 0.2,
          row_y + ch / 2, color=BORDER_STR, weight=1.0)

    for i, (label, accent) in enumerate(steps):
        x = pad_x + i * (cw + gap)
        # Card
        round_rect(s, x, row_y, cw, ch, fill=WHITE, line=BORDER, radius=0.10)
        # Top accent stripe
        rect(s, x, row_y, cw, 0.10, fill=accent)
        # Number circle in top-center
        d = 0.50
        numbered_circle(s, x + cw / 2 - d / 2, row_y + 0.30, d, i + 1,
                        bg=accent, fg=WHITE, size=14)
        # Label
        text(s, x + 0.10, row_y + 0.95, cw - 0.20, 1.10,
             label, size=12, color=NAVY_900, bold=True,
             align="center", valign="middle")

    # Roles strip below the timeline — who acts at each step
    strip_y = 5.45
    strip_h = 0.40
    rect(s, pad_x, strip_y, avail, strip_h, fill=BG_SOFT)
    text(s, pad_x + 0.30, strip_y, 2.8, strip_h,
         "¿Quién actúa?", size=11, color=TEXT_MUTED, bold=True,
         valign="middle", charspace=120)
    actor_labels = [
        ("ASESOR",   NAVY_900,  3.50),
        ("IA",       CYAN_500,  5.20),
        ("ASESOR",   EMERALD,   6.40),
        ("SISTEMA",  AMBER,     7.70),
        ("ASESOR",   EMERALD,   9.70),
        ("ASESOR",   EMERALD,  10.65),
        ("COMPLIANCE", VIOLET, 11.80),
    ]
    # Simplified: just say "AI / Advisor / System / Compliance" markers under groups
    # Using a single explanation line
    text(s, pad_x, strip_y + 0.45, avail, 0.32,
         "Asesor + IA preparan el perfil  →  Asesor aprueba  →  Sistema controla suitability y riesgo  →  Asesor selecciona  →  Compliance verifica.",
         size=11, color=TEXT_MUTED, italic=True, align="center")

    # Pull-quote at bottom
    text(s, pad_x, 6.45, avail, 0.45,
         "“Si alguno de los nueve pasos no se firmó, el reporte no se genera.”",
         size=14, color=NAVY_900, italic=True, align="center", bold=True)

    footer(s, 4)


# ─────────────────────────────────────────────────────────────────────
# Slide 5 — Diferencial risk-first
# ─────────────────────────────────────────────────────────────────────
def slide_5_diferencial():
    s = add_slide(bg=BG_LIGHT)

    eyebrow_chip(s, 0.85, 0.6, "DIFERENCIAL · RISK-FIRST",
                 fg=WHITE, bg=NAVY_900, w=3.0)
    text(s, 0.85, 1.05, 11.6, 0.9,
         "Empezamos por el riesgo aceptable, no por el retorno deseado",
         size=28, color=NAVY_900, bold=True)
    text(s, 0.85, 1.85, 11.6, 0.45,
         "“Risk-first, not return-chasing.”",
         size=15, color=TEXT_MUTED, italic=True)

    # Two-column comparison
    col_w = 5.7
    col_h = 2.55
    cy = 2.55
    gap = 0.30
    left_x = 0.85
    right_x = left_x + col_w + gap

    # LEFT — Tradicional
    card(s, left_x, cy, col_w, col_h, fill=WHITE, accent=ROSE, accent_w=0.08)
    text(s, left_x + 0.35, cy + 0.25, col_w - 0.5, 0.35,
         "ENFOQUE TRADICIONAL · RETURN-CHASING",
         size=10.5, color=ROSE, bold=True, charspace=160)
    text(s, left_x + 0.35, cy + 0.75, col_w - 0.5, 0.45,
         "Empieza por la expectativa de retorno", size=15,
         color=NAVY_900, bold=True)
    text_lines(s, left_x + 0.35, cy + 1.25, col_w - 0.5, 1.25, [
        "•  Cliente pide “10% anual en dólares”.",
        "•  Asesor busca productos que pinten ese número.",
        "•  El riesgo aparece como consecuencia, no como límite.",
        "•  El override de límites no queda firmado por escrito.",
    ], size=12, color=TEXT_BODY, line_spacing=1.25)

    # RIGHT — Risk-first
    card(s, right_x, cy, col_w, col_h, fill=WHITE, accent=EMERALD, accent_w=0.08)
    text(s, right_x + 0.35, cy + 0.25, col_w - 0.5, 0.35,
         "RISK-FIRST ADVISORY",
         size=10.5, color=EMERALD, bold=True, charspace=160)
    text(s, right_x + 0.35, cy + 0.75, col_w - 0.5, 0.45,
         "Empieza por la tolerancia y capacidad de riesgo", size=15,
         color=NAVY_900, bold=True)
    text_lines(s, right_x + 0.35, cy + 1.25, col_w - 0.5, 1.25, [
        "•  KYC fija el presupuesto de riesgo aprobado.",
        "•  El optimizador busca el mejor retorno dentro de ese límite.",
        "•  Exceder el límite requiere override firmado por el asesor.",
        "•  Todo queda atado a la cadena de auditoría del caso.",
    ], size=12, color=TEXT_BODY, line_spacing=1.25)

    # Three variants strip
    sy = 5.4
    sh = 1.55
    text(s, 0.85, sy - 0.45, 11.6, 0.35,
         "El optimizador genera tres variantes — y bloquea las que se pasan del presupuesto sin override:",
         size=12, color=TEXT_MUTED)

    variants = [
        ("DEFENSIVE", "Por debajo del presupuesto aprobado",  NAVY_700,  False),
        ("BALANCED",  "Dentro del presupuesto aprobado",       EMERALD,   False),
        ("GROWTH",    "Puede exceder el presupuesto · requiere override firmado", AMBER, True),
    ]
    vw = (11.6 - 2 * 0.25) / 3
    for i, (name, body, color_, needs_override) in enumerate(variants):
        x = 0.85 + i * (vw + 0.25)
        round_rect(s, x, sy, vw, sh, fill=WHITE, line=BORDER, radius=0.08)
        rect(s, x, sy, vw, 0.10, fill=color_)
        text(s, x + 0.35, sy + 0.30, vw - 0.5, 0.45, name,
             size=20, color=NAVY_900, bold=True)
        text(s, x + 0.35, sy + 0.85, vw - 0.5, 0.7, body,
             size=11.5, color=TEXT_BODY)
        if needs_override:
            round_rect(s, x + vw - 1.75, sy + 0.30, 1.40, 0.32,
                       fill=AMBER_LT, line=AMBER, radius=0.4)
            text(s, x + vw - 1.75, sy + 0.30, 1.40, 0.32,
                 "OVERRIDE", size=9, color=AMBER, bold=True,
                 align="center", valign="middle", charspace=200)

    footer(s, 5)


# ─────────────────────────────────────────────────────────────────────
# Slide 6 — Qué ve el asesor en la demo
# ─────────────────────────────────────────────────────────────────────
def slide_6_demo():
    s = add_slide(bg=BG_LIGHT)

    eyebrow_chip(s, 0.85, 0.6, "QUÉ VE EL ASESOR EN LA DEMO",
                 fg=WHITE, bg=NAVY_900, w=3.4)
    text(s, 0.85, 1.05, 11.6, 0.9,
         "Una demo en español, paso a paso, que un asesor entiende sola",
         size=26, color=NAVY_900, bold=True)
    text(s, 0.85, 1.85, 11.6, 0.45,
         "Pensada para mostrar el producto a profesores de finanzas o asesores no técnicos en ~3 minutos.",
         size=14, color=TEXT_MUTED)

    features = [
        ("Demo en español",
         "Hero, formulario, botones y mensajes localizados. El asesor no necesita traducir mentalmente nada."),
        ("Perfil inversor paso a paso",
         "8 cards didácticas: cada paso explica qué hace, por qué importa y qué mirar en el resultado."),
        ("Propuestas con instrumentos",
         "Cada variante muestra los tickers, pesos, tipo, moneda y motivo de inclusión — no un score abstracto."),
        ("Reporte Markdown",
         "Vista previa del documento que se entregaría al cliente, con composición y comparación de variantes."),
        ("Auditoría intacta",
         "Pill verde “cadena intacta” sobre los eventos del caso. Compliance ve la integridad de un vistazo."),
        ("AI logs con PII redactada",
         "Cada llamada a OpenAI guardada con el texto libre y el client_id redactados por el backend."),
    ]
    cols = 3
    rows = 2
    pad_x = 0.85
    pad_top = 2.55
    gap = 0.25
    cw = (SW - 2 * pad_x - (cols - 1) * gap) / cols
    ch = 2.05

    for i, (title, body) in enumerate(features):
        col = i % cols
        row = i // cols
        x = pad_x + col * (cw + gap)
        y = pad_top + row * (ch + gap)
        card(s, x, y, cw, ch, fill=WHITE, accent=CYAN_500, accent_w=0.08)
        # Number circle
        numbered_circle(s, x + 0.35, y + 0.30, 0.45, i + 1,
                        bg=CYAN_500, fg=WHITE, size=13)
        text(s, x + 1.00, y + 0.27, cw - 1.2, 0.45, title,
             size=14, color=NAVY_900, bold=True)
        text(s, x + 0.35, y + 0.90, cw - 0.55, ch - 1.05, body,
             size=11.5, color=TEXT_BODY)

    footer(s, 6)


# ─────────────────────────────────────────────────────────────────────
# Slide 7 — Propuestas de cartera (holdings table)
# ─────────────────────────────────────────────────────────────────────
def slide_7_propuestas():
    s = add_slide(bg=BG_LIGHT)

    eyebrow_chip(s, 0.85, 0.6, "PROPUESTAS DE CARTERA",
                 fg=WHITE, bg=NAVY_900, w=2.8)
    text(s, 0.85, 1.05, 11.6, 0.9,
         "Cada variante muestra instrumentos, pesos y motivo",
         size=28, color=NAVY_900, bold=True)
    text(s, 0.85, 1.85, 11.6, 0.45,
         "“Portfolio proposals with instruments, weights and rationale.”",
         size=14, color=TEXT_MUTED, italic=True)

    # LEFT: variant card (mockup of BALANCED) with holdings table
    lx = 0.85
    ly = 2.50
    lw = 7.7
    lh = 4.45
    round_rect(s, lx, ly, lw, lh, fill=WHITE, line=BORDER, radius=0.08)
    rect(s, lx, ly, lw, 0.10, fill=EMERALD)

    text(s, lx + 0.30, ly + 0.25, lw - 0.60, 0.45,
         "Variante BALANCED  —  dentro del presupuesto de riesgo",
         size=14, color=NAVY_900, bold=True)
    text(s, lx + 0.30, ly + 0.70, lw - 0.60, 0.30,
         "Retorno esperado 6,8 %  ·  Volatilidad 9,4 %  ·  6 instrumentos",
         size=11, color=TEXT_MUTED)

    # Table headers
    hy = ly + 1.10
    rh = 0.42
    col_x = [lx + 0.30, lx + 1.85, lx + 3.30, lx + 4.55, lx + 5.20, lx + 5.95]
    col_w_arr = [1.55, 1.45, 1.25, 0.65, 0.75, 1.55]
    headers = ["INSTRUMENTO", "TIPO", "MONEDA", "PESO", "BARRA", "MOTIVO"]
    rect(s, lx + 0.20, hy, lw - 0.40, rh, fill=BG_SOFT)
    for i, h in enumerate(headers):
        if h == "BARRA":
            continue  # no header for the bar column
        align = "right" if h == "PESO" else "left"
        text(s, col_x[i], hy, col_w_arr[i], rh, h,
             size=9, color=TEXT_MUTED, bold=True,
             valign="middle", align=align, charspace=120)

    # Rows  — sample holdings, illustrative
    rows = [
        ("GD30",   "Bono soberano",      "USD", 0.30, "Renta fija · hard dollar"),
        ("AL30",   "Bono soberano",      "USD", 0.22, "Renta fija · hard dollar"),
        ("YPFD",   "Acción",             "ARS", 0.15, "Equity · energía"),
        ("SHV",    "ETF money market",   "USD", 0.13, "Liquidez · USD"),
        ("PAMP",   "CEDEAR",             "USD", 0.10, "Equity · diversificación"),
        ("MEP",    "Money market USD",   "USD", 0.10, "Liquidez · USD"),
    ]
    max_w = max(w for _, _, _, w, _ in rows)
    for i, (ticker, kind, cur, weight, motivo) in enumerate(rows):
        ry = hy + rh + 0.05 + i * (rh + 0.05)
        rect(s, lx + 0.20, ry, lw - 0.40, rh, fill=WHITE)
        if i < len(rows) - 1:
            hline(s, lx + 0.30, ry + rh, lx + lw - 0.30,
                  ry + rh, color=BORDER, weight=0.4)
        text(s, col_x[0], ry, col_w_arr[0], rh, ticker,
             size=12, color=NAVY_900, bold=True, font=FONT_MONO,
             valign="middle")
        text(s, col_x[1], ry, col_w_arr[1], rh, kind,
             size=11, color=TEXT_BODY, valign="middle")
        text(s, col_x[2], ry, col_w_arr[2], rh, cur,
             size=11, color=TEXT_BODY, valign="middle")
        text(s, col_x[3], ry, col_w_arr[3], rh, f"{int(weight*100)} %",
             size=11.5, color=NAVY_900, bold=True,
             valign="middle", align="right")
        # weight bar
        bar_x = col_x[4]
        bar_w = col_w_arr[4]
        bar_y = ry + rh / 2 - 0.07
        rect(s, bar_x, bar_y, bar_w, 0.14, fill=BG_SOFT)
        fill_w = bar_w * (weight / max_w)
        rect(s, bar_x, bar_y, fill_w, 0.14, fill=NAVY_600)
        # motivo
        text(s, col_x[5], ry, col_w_arr[5], rh, motivo,
             size=10.5, color=TEXT_MUTED, valign="middle", italic=True)

    # RIGHT: what each candidate exposes
    rx = lx + lw + 0.30
    rw = SW - rx - 0.85
    text(s, rx, 2.50, rw, 0.40,
         "QUÉ EXPONE CADA VARIANTE",
         size=10.5, color=NAVY_600, bold=True, charspace=160)
    text(s, rx, 2.95, rw, 0.40,
         "Para cada candidato propuesto",
         size=16, color=NAVY_900, bold=True)
    items = [
        ("Instrumentos",  "ticker + nombre"),
        ("Pesos",         "decimal y porcentaje"),
        ("Tipo",          "bono / ETF / equity / CEDEAR"),
        ("Moneda",        "USD / ARS / EUR"),
        ("Motivos",       "asset class, hard dollar, etc."),
        ("Override",      "si excede el presupuesto"),
    ]
    iy = 3.55
    for i, (k, v) in enumerate(items):
        oy = iy + i * 0.55
        oval(s, rx, oy + 0.10, 0.16, fill=CYAN_500)
        text(s, rx + 0.30, oy, rw - 0.30, 0.40, k,
             size=12, color=NAVY_900, bold=True)
        text(s, rx + 0.30, oy + 0.22, rw - 0.30, 0.30, v,
             size=10.5, color=TEXT_MUTED, italic=True)

    footer(s, 7)


# ─────────────────────────────────────────────────────────────────────
# Slide 8 — Compliance y auditoría
# ─────────────────────────────────────────────────────────────────────
def slide_8_compliance():
    s = add_slide(bg=BG_LIGHT)

    eyebrow_chip(s, 0.85, 0.6, "COMPLIANCE Y AUDITORÍA",
                 fg=WHITE, bg=NAVY_900, w=2.9)
    text(s, 0.85, 1.05, 11.6, 0.9,
         "Trazabilidad por diseño, no agregada al final",
         size=28, color=NAVY_900, bold=True)
    text(s, 0.85, 1.85, 11.6, 0.45,
         "“Every decision auditable.”",
         size=15, color=TEXT_MUTED, italic=True)

    items = [
        ("Hash chain SHA-256",
         "Cada decisión (KYC, análisis IA, aprobación, override, selección, reporte) queda en una cadena SHA-256 por caso. Si alguien edita un payload, la cadena se rompe.",
         VIOLET, "01"),
        ("AI logs con PII redactada",
         "Cada llamada a OpenAI se guarda con el texto libre y el client_id redactados por el backend antes de persistir. No queda PII del cliente en el log.",
         CYAN_500, "02"),
        ("Decisiones humanas firmadas",
         "El asesor aprueba o rechaza con un rationale obligatorio. El override del presupuesto de riesgo requiere firma escrita para poder seleccionar la variante.",
         EMERALD, "03"),
        ("Reportes revisables",
         "El reporte final es Markdown determinístico. El asesor lo revisa antes de compartir con el cliente. No hay botón de “enviar al cliente”.",
         NAVY_600, "04"),
    ]
    cw = (11.6 - 0.30) / 2
    ch = 2.20
    cx0 = 0.85
    cy0 = 2.55

    for i, (title, body, accent, num) in enumerate(items):
        col = i % 2
        row = i // 2
        x = cx0 + col * (cw + 0.30)
        y = cy0 + row * (ch + 0.20)
        card(s, x, y, cw, ch, fill=WHITE, accent=accent, accent_w=0.08)
        text(s, x + 0.35, y + 0.30, 0.50, 0.45, num,
             size=22, color=accent, bold=True, font=FONT_MONO)
        text(s, x + 1.00, y + 0.30, cw - 1.20, 0.45, title,
             size=15, color=NAVY_900, bold=True)
        text(s, x + 0.35, y + 0.95, cw - 0.55, ch - 1.10, body,
             size=11.5, color=TEXT_BODY)

    # Bottom callout
    text(s, 0.85, 7.05, 11.6, 0.35,
         "El sistema no es una “recomendación automática” — es un copilot que prepara el material para que el asesor decida.",
         size=12, color=NAVY_900, italic=True, bold=True, align="center")

    footer(s, 8)


# ─────────────────────────────────────────────────────────────────────
# Slide 9 — Estado actual del producto
# ─────────────────────────────────────────────────────────────────────
def slide_9_estado_actual():
    s = add_slide(bg=BG_LIGHT)

    eyebrow_chip(s, 0.85, 0.6, "ESTADO ACTUAL", fg=WHITE, bg=EMERALD, w=1.7)
    text(s, 0.85, 1.05, 11.6, 0.9,
         "Lo que ya funciona, hoy, en una máquina local",
         size=30, color=NAVY_900, bold=True)
    text(s, 0.85, 1.85, 11.6, 0.45,
         "Un prototipo funcional end-to-end, no un mockup. Operable en ~3 minutos tras correr el bootstrap.",
         size=14, color=TEXT_MUTED)

    # LEFT — checklist
    lx = 0.85
    ly = 2.55
    lw = 7.7
    lh = 4.45
    round_rect(s, lx, ly, lw, lh, fill=WHITE, line=BORDER, radius=0.08)
    text(s, lx + 0.40, ly + 0.30, lw - 0.5, 0.4,
         "ENTREGADO", size=10.5, color=EMERALD, bold=True, charspace=160)

    deliverables = [
        "Backend case-scoped completo (Fase 2)",
        "Dashboard técnico de firms / asesores / clientes / casos",
        "Workbench de 15 paneles end-to-end",
        "Demo investor-facing en español, paso a paso",
        "Bootstrap local en un comando",
        "Seed de datos demo (firm + asesor + cliente + caso)",
        "Composición de carteras con instrumentos y pesos",
        "Hash chain por caso + AI logs con PII redactada",
    ]
    for i, d in enumerate(deliverables):
        iy = ly + 0.80 + i * 0.45
        # Check icon: small emerald circle + white check (use a triangle as a proxy of a check)
        oval(s, lx + 0.40, iy + 0.08, 0.26, fill=EMERALD)
        text(s, lx + 0.40, iy + 0.08, 0.26, 0.26, "✓",
             size=12, color=WHITE, bold=True, align="center", valign="middle")
        text(s, lx + 0.85, iy, lw - 1.0, 0.40, d,
             size=12.5, color=TEXT_DARK, valign="middle")

    # RIGHT — Big stat
    rx = lx + lw + 0.30
    rw = SW - rx - 0.85
    round_rect(s, rx, ly, rw, 2.05, fill=NAVY_900, radius=0.08)
    text(s, rx + 0.35, ly + 0.30, rw - 0.5, 0.4,
         "SUITE DE TESTS", size=11, color=CYAN_LIGHT, bold=True, charspace=160)
    text(s, rx + 0.35, ly + 0.65, rw - 0.5, 1.0,
         "3 087", size=70, color=WHITE, bold=True)
    text(s, rx + 0.35, ly + 1.65, rw - 0.5, 0.4,
         "tests verdes  ·  unit + integration",
         size=12, color=SOFT_INK)

    # Below stat — three context tiles (compact, fit within slide bottom)
    tile_y = ly + 2.20
    tile_h = 0.70
    tile_gap = 0.05
    items = [
        ("100 %",   "Casos backend cubiertos"),
        ("0",       "Dependencias frontend"),
        ("~3 min",  "De cero a demo operable"),
    ]
    for i, (val, label) in enumerate(items):
        ty = tile_y + i * (tile_h + tile_gap)
        round_rect(s, rx, ty, rw, tile_h, fill=WHITE, line=BORDER, radius=0.06)
        text(s, rx + 0.25, ty + 0.05, rw - 0.4, 0.36, val,
             size=18, color=NAVY_900, bold=True)
        text(s, rx + 0.25, ty + 0.38, rw - 0.4, 0.28, label,
             size=10, color=TEXT_MUTED)

    footer(s, 9)


# ─────────────────────────────────────────────────────────────────────
# Slide 10 — Limitaciones
# ─────────────────────────────────────────────────────────────────────
def slide_10_limitaciones():
    s = add_slide(bg=BG_LIGHT)

    eyebrow_chip(s, 0.85, 0.6, "LIMITACIONES ACTUALES", fg=WHITE, bg=AMBER, w=2.7)
    text(s, 0.85, 1.05, 11.6, 0.9,
         "Honestidad sobre el alcance: esto NO es production-ready",
         size=28, color=NAVY_900, bold=True)
    text(s, 0.85, 1.85, 11.6, 0.45,
         "El prototipo funcional muestra el flujo end-to-end. Para piloto productivo, hay que cerrar todo lo siguiente.",
         size=14, color=TEXT_MUTED)

    limits = [
        ("No production-ready",
         "Demo local · pilot interno con datos demo. No usar con PII real de clientes."),
        ("SQLite local",
         "Sin PostgreSQL, sin réplica, sin clustering. Suficiente para demo, no para piloto."),
        ("Auth con dev tokens",
         "Bearer tokens opacos en YAML. Sin JWT, sin IdP, sin rotación, sin revocation."),
        ("Universo de instrumentos demo",
         "20 instrumentos del fixture CSV. No es market data productivo ni tiene SLA de frescura."),
        ("Sin market data provider real",
         "Los snapshots vienen del fixture. Falta integración con Bloomberg / Refinitiv / proveedor real."),
        ("Reporte solo Markdown",
         "Sin PDF, sin branding de la firm. La impresión / entrega al cliente sigue siendo manual."),
        ("Sin firm-level access control completo",
         "Cualquier token con rol válido puede ver cualquier caso. No hay aislación entre firms."),
        ("Sin datos reales sensibles",
         "Recordatorio: no cargar PII real ni información confidencial en esta demo."),
    ]
    cols = 2
    cx0 = 0.85
    cy0 = 2.55
    cw = (11.6 - 0.30) / cols
    ch = 1.05
    gap_y = 0.12

    for i, (title, body) in enumerate(limits):
        col = i % cols
        row = i // cols
        x = cx0 + col * (cw + 0.30)
        y = cy0 + row * (ch + gap_y)
        card(s, x, y, cw, ch, fill=WHITE, accent=AMBER, accent_w=0.06)
        text(s, x + 0.30, y + 0.15, cw - 0.45, 0.36, title,
             size=13, color=NAVY_900, bold=True)
        text(s, x + 0.30, y + 0.50, cw - 0.45, 0.55, body,
             size=10.5, color=TEXT_BODY)

    footer(s, 10)


# ─────────────────────────────────────────────────────────────────────
# Slide 11 — Roadmap (Fase 4)
# ─────────────────────────────────────────────────────────────────────
def slide_11_roadmap():
    s = add_slide(bg=BG_LIGHT)

    eyebrow_chip(s, 0.85, 0.6, "ROADMAP · FASE 4", fg=WHITE, bg=CYAN_500, w=2.0)
    text(s, 0.85, 1.05, 11.6, 0.9,
         "De prototipo funcional a pilot readiness",
         size=30, color=NAVY_900, bold=True)
    text(s, 0.85, 1.85, 11.6, 0.45,
         "Cierre del prototipo + hardening necesario para correr la demo con un asesor piloto real.",
         size=14, color=TEXT_MUTED)

    items = [
        ("Market data provider productivo",
         "Reemplazar el CSV fixture por un live provider con SLA de frescura y validación."),
        ("Manual universe upload",
         "Admin endpoint para reemplazar el universo sin redeploy."),
        ("PDF / branding del reporte",
         "Render del Markdown a PDF con logo + colores + header de la firm."),
        ("Firm-level access control completo",
         "firm_id en el token + filtrado por firm en todos los endpoints /cases/*."),
        ("Auth productiva (JWT / OIDC)",
         "Integración con identity provider. Rotación, revocation, emisión de tokens."),
        ("Backup / restore",
         "Snapshot regular de la DB SQLite (o migración a PostgreSQL si la escala lo pide)."),
        ("Pilot readiness checklist",
         "Cifrado at-rest, retention policy, sign-off legal, runbook de incidentes."),
    ]
    n = len(items)
    # Single column with cyan timeline on the left
    tx = 1.10
    timeline_x = tx + 0.22 / 2
    cy0 = 2.55
    h_per = (7.0 - cy0) / n

    # Vertical timeline line
    hline(s, timeline_x, cy0 + 0.20, timeline_x,
          cy0 + (n - 1) * h_per + 0.30, color=BORDER_STR, weight=2.0)

    for i, (title, body) in enumerate(items):
        y = cy0 + i * h_per
        # Numbered dot
        oval(s, tx, y + 0.15, 0.34, fill=CYAN_500)
        text(s, tx, y + 0.15, 0.34, 0.34, str(i + 1),
             size=12, color=WHITE, bold=True, align="center", valign="middle")
        # Card
        cx = tx + 0.70
        cw = SW - cx - 0.85
        text(s, cx, y + 0.12, cw, 0.36, title,
             size=14, color=NAVY_900, bold=True)
        text(s, cx, y + 0.45, cw, 0.35, body,
             size=11, color=TEXT_BODY)

    footer(s, 11)


# ─────────────────────────────────────────────────────────────────────
# Slide 12 — Cierre / pitch
# ─────────────────────────────────────────────────────────────────────
def slide_12_cierre():
    s = add_slide(bg=NAVY_DARK)

    # Decorative thin lines on the left edge
    for x in [0.50, 0.70, 0.90]:
        hline(s, x, 0.0, x, SH, color=NAVY_700, weight=0.5)

    # Brand mark at the top
    oval(s, 1.30, 0.78, 0.20, fill=CYAN_500)
    text(s, 1.63, 0.66, 6.0, 0.5,
         "RISK-FIRST  ADVISORY",
         size=12, color=SKY_LIGHT, bold=True, charspace=240)

    # Big quote
    text_lines(s, 1.30, 2.10, 11.0, 3.0, [
        [("La asesoría asistida por IA debería ser ",
          {"size": 32, "color": LIGHT_INK})],
        [("controlada, justificable y auditable.",
          {"size": 32, "color": WHITE, "bold": True})],
        [("La tecnología potencia al asesor ",
          {"size": 32, "color": LIGHT_INK})],
        [("sin reemplazarlo.",
          {"size": 32, "color": CYAN_500, "bold": True})],
    ], line_spacing=1.20, space_after=4)

    # English tagline cluster
    text_lines(s, 1.30, 5.55, 11.0, 1.0, [
        [("“AI proposes. Advisor decides.  ", {"size": 16, "color": CYAN_LIGHT, "italic": True}),
         ("Risk-first, not return-chasing.  ", {"size": 16, "color": CYAN_LIGHT, "italic": True}),
         ("Every decision auditable.”",        {"size": 16, "color": CYAN_LIGHT, "italic": True})],
    ], line_spacing=1.2)

    # Bottom info row
    by = 6.50
    bh = 0.7
    bw = (SW - 2.6) / 3
    items = [
        ("ESTADO",       "Prototipo funcional · demo local"),
        ("PRÓXIMO HITO", "Fase 4 — pilot readiness"),
        ("CONTACTO",     "Pedí la demo · ~3 min en vivo"),
    ]
    for i, (label, val) in enumerate(items):
        x = 1.30 + i * (bw + 0.20)
        round_rect(s, x, by, bw, bh, fill=NAVY_800, line=NAVY_700, radius=0.10)
        text(s, x + 0.25, by + 0.10, bw - 0.5, 0.32,
             label, size=10, color=CYAN_LIGHT, bold=True, charspace=160)
        text(s, x + 0.25, by + 0.36, bw - 0.5, 0.32,
             val, size=12.5, color=WHITE, bold=True)


# ─────────────────────────────────────────────────────────────────────
# Build
# ─────────────────────────────────────────────────────────────────────
def main():
    slide_1_portada()
    slide_2_problema()
    slide_3_idea_central()
    slide_4_flujo()
    slide_5_diferencial()
    slide_6_demo()
    slide_7_propuestas()
    slide_8_compliance()
    slide_9_estado_actual()
    slide_10_limitaciones()
    slide_11_roadmap()
    slide_12_cierre()

    out = "Risk-First-Advisory-Pitch-Deck.pptx"
    prs.save(out)
    print(f"Wrote {out} ({len(prs.slides)} slides, {prs.slide_width.inches}x{prs.slide_height.inches} in)")


if __name__ == "__main__":
    main()
